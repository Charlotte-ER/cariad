import os
import pandas as pd
import sys
import shutil
import unittest

sys.path.append('../src/cariad')
from cariad import run, parse_cli_args

from pptx import Presentation
from docx import Document

class Test(unittest.TestCase):

    @classmethod
    def setUpClass(cls):
        cls.case_reference = 'CAS-12345-EXAMPLE'

        # Set up example folder
        os.mkdir(f'./{cls.case_reference}')

        # Create example dataset
        cls.data = {
            'Print Index Number': [1, 2, 3, 4, 5, 6, 7, 8, 9], 
            'Search': ['A', 'A', 'A', 'A', 'A', 'A', 'A', 'A', 'A'],
            'Database': ['DB_NAME', 'DB_NAME', 'DB_NAME', 'DB_NAME', 'DB_NAME', 'DB_NAME', 'DB_NAME', 'DB_NAME', 'DB_NAME'],
            'Doc. Number': ['1234561', '1234562', '1234567', '1234568', '1234569', '1234570', '1234571', '1234572', '1234573'],
            'Version': [1, 1, 1, 1, 1, 1, 1, 1, 1],
            'Extension': ['PPT', 'DOC', 'PDF', 'PPTX', 'DOCX', 'DOCX', 'MSG', 'MSG', 'MSG']
        }
        df = pd.DataFrame(cls.data)

        # Create example spreadsheet
        test_spreadsheet_path = f'{cls.case_reference}/{cls.case_reference}.xlsx'
        df.to_excel(test_spreadsheet_path, sheet_name='Index', header=True, index=False, startrow=2)

        # Create example documents
        for _, row in df.iterrows():
            filename = f"./{cls.case_reference}/{row['Doc. Number']}_{row['Version']}.{row['Extension']}"

            if row['Extension'] == 'PPT' or row['Extension'] == 'PPTX':
                presentation = Presentation()
                layout = presentation.slide_masters[0].slide_layouts[0]
                presentation.slides.add_slide(layout)
                presentation.save(filename)

            elif row['Extension'] == 'DOC' or row['Extension'] == 'DOCX':
                document = Document()
                document.save(filename)  


    @classmethod
    def tearDownClass(cls):
        # Delete example folder and contents
        shutil.rmtree(f'./{cls.case_reference}')

    def test_no_args(self):
        sys.argv = ['cariad.py']
        with self.assertRaises(SystemExit):
            parse_cli_args()


    def test_only_one_arg(self):
        sys.argv = ['cariad.py', 'CAS-12345-EXAMPLE.xlsx']
        with self.assertRaises(SystemExit):
            parse_cli_args()


    def test_args_in_wrong_order(self):
        sys.argv = ['cariad.py', 'CAS-12345-EXAMPLE.xlsx', './CAS-12345-EXAMPLE']
        with self.assertRaises(SystemExit):
            parse_cli_args()


    def test_unknown_arguments(self):
        os.makedirs(f'./Two Words')
        shutil.copyfile('./CAS-12345-EXAMPLE/CAS-12345-EXAMPLE.xlsx', './Two Words/CAS-12345-EXAMPLE.xlsx')
        sys.argv = ['cariad.py', './Two', 'Words', 'CAS-12345-EXAMPLE.xlsx']
        self.assertEqual(parse_cli_args(), ('./Two Words', './Two Words\\CAS-12345-EXAMPLE.xlsx'))
        shutil.rmtree(f'./Two Words')


    def test_directory_exists(self):
        # Directory exists
        sys.argv = ['cariad.py', './CAS-12345-EXAMPLE', 'CAS-12345-EXAMPLE.xlsx']
        self.assertEqual(parse_cli_args(), ('./CAS-12345-EXAMPLE', './CAS-12345-EXAMPLE\\CAS-12345-EXAMPLE.xlsx'))


        # Directory does not exist
        sys.argv = ['cariad.py', './CAS-BAD-EXAMPLE', 'CAS-12345-EXAMPLE.xlsx']
        with self.assertRaises(SystemExit):
            parse_cli_args()
    

    def test_spreadsheet_exists(self):
        sys.argv = ['cariad.py', './CAS-12345-EXAMPLE', 'CAS-BAD-EXAMPLE.xlsx']
        with self.assertRaises(SystemExit):
            parse_cli_args()


    def test_reference_is_a_spreadsheet(self):
        with self.assertRaises(SystemExit):
            run.get_index_dataframe_from_spreadsheet('./CAS-12345-EXAMPLE/CAS-12345-EXAMPLE.docx', 'cli')


    def test_spreadsheet_has_index_tab(self):
        # Has an Index tab
        self.assertEqual(isinstance(run.get_index_dataframe_from_spreadsheet('./CAS-12345-EXAMPLE/CAS-12345-EXAMPLE.xlsx', 'cli'), 
                                    pd.DataFrame), True)
        # Does NOT have an Index tab - TODO


if __name__ == '__main__':
    unittest.main()